import re
import typing

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl.cell.cell import Cell


class PPTXStyle:
    def __init__(self, pptx_file_path: str):
        self.pptx_file_path = pptx_file_path
        self.pptx_obj = Presentation(pptx_file_path)

    def replace_text(
            self,
            slide_index: int,
            old_text_pattern: str,
            new_text_parts: list,
            # use_regex: bool = False
    ):
        """
        替换指定幻灯片中的文本。

        :param slide_index: 幻灯片索引（从0开始）
        :param old_text_pattern: 要替换的旧文本（或正则表达式）
        :param new_text_parts: 为了保留格式，必须逐个块进行编辑操作
        """
        slide = self.pptx_obj.slides[slide_index]
        print(f"slide_index:{slide_index}", end="")
        self.__recursive_replace_text_in_shapes(slide.shapes, old_text_pattern, new_text_parts)
        return self

    def __recursive_replace_text_in_shapes(self, shapes, old_text_pattern: str, new_text_parts: list):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.__recursive_replace_text_in_shapes(shape.shapes, old_text_pattern, new_text_parts)
                continue
            elif not shape.has_text_frame:
                continue
            # 正常路径
            for paragraph in shape.text_frame.paragraphs:
                if re.search(old_text_pattern, paragraph.text):
                    # 如果不写 new_text_parts，那么输出当前所有run
                    if not new_text_parts:
                        r = ['"' + run.text + '"' for run in paragraph.runs]
                        print("["+",".join(r),"]")
                        return
                    for run, new_text_part in zip(paragraph.runs, new_text_parts):
                        if new_text_part is None:
                            continue
                        # 使用正则查找并替换
                        if isinstance(new_text_part, str):
                            run.text = new_text_part
                        elif isinstance(new_text_part, typing.Callable):
                            new_text = new_text_part(run.text)
                            run.text = new_text

    def replace_pic(
            self,
            slide_index: int,
            old_pic_id,
            new_pic_path: str,
    ):
        """
        替换指定幻灯片中的图片。

        :param slide_index: 幻灯片索引（从0开始）
        :param old_pic_id: 用于识别旧图片的标识。
                           - 若为 int：匹配 shape.shape_id
                           - 若为 str：匹配 shape.name
        :param new_pic_path: 新图片的本地文件路径
        """
        slide = self.pptx_obj.slides[slide_index]
        print(f"slide_index:{slide_index}", end="")

        # 遍历所有 shapes（包括组合内的）
        def find_picture_in_shapes(shapes):
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    result = find_picture_in_shapes(shape.shapes)
                    if result is not None:
                        return result
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    print(f"pic: shape_id: {shape.shape_id}, shape_name: {shape.name}")
                    if isinstance(old_pic_id, int) and shape.shape_id == old_pic_id:
                        return shape
                    elif isinstance(old_pic_id, str) and shape.name == old_pic_id:
                        return shape
            return None

        target_shape = find_picture_in_shapes(slide.shapes)

        if target_shape is None:
            raise ValueError(f"未找到匹配的图片 (old_pic_id={old_pic_id})")

        # 记录原始位置和尺寸
        left = target_shape.left
        top = target_shape.top
        width = target_shape.width
        height = target_shape.height

        # 删除旧图片
        sp = target_shape._element
        sp.getparent().remove(sp)

        # 插入新图片（保持原位置和尺寸）
        slide.shapes.add_picture(
            new_pic_path,
            left=left,
            top=top,
            width=width,
            height=height
        )

        return self

    def replace_table(
            self,
            slide_index: int,
            old_table_id,
            new_table_excel: str,
            new_table_sheet_id,
            new_table_value_range: str,
    ):
        """
        替换指定幻灯片中的表格内容，从 Excel 指定区域读取数据并**保留显示格式**。
        """
        slide = self.pptx_obj.slides[slide_index]
        target_table_shape = None

        def find_table_in_shapes(shapes):
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    result = find_table_in_shapes(shape.shapes)
                    if result is not None:
                        return result
                elif shape.has_table:
                    if isinstance(old_table_id, int) and shape.shape_id == old_table_id:
                        return shape
                    elif isinstance(old_table_id, str) and shape.name == old_table_id:
                        return shape
            return None

        target_table_shape = find_table_in_shapes(slide.shapes)
        if target_table_shape is None:
            raise ValueError(f"未找到匹配的表格 (old_table_id={old_table_id})")

        table = target_table_shape.table

        # 读取 Excel（注意：必须关闭 read_only 才能获取 number_format！）
        from openpyxl import load_workbook
        wb = load_workbook(new_table_excel, read_only=False, data_only=True)
        try:
            if isinstance(new_table_sheet_id, int):
                ws = wb.worksheets[new_table_sheet_id]
            else:
                ws = wb[new_table_sheet_id]
        except (KeyError, IndexError) as e:
            wb.close()
            raise ValueError(f"Excel 中未找到指定的 sheet: {new_table_sheet_id}") from e

        try:
            excel_range = ws[new_table_value_range]
        except Exception as e:
            wb.close()
            raise ValueError(f"无效的 Excel 区域: {new_table_value_range}") from e

        # 提取格式化后的文本（二维列表）
        new_display_texts = []
        for row in excel_range:
            new_display_texts.append([self._format_cell_value(cell) for cell in row])

        wb.close()

        # 检查行列匹配
        ppt_rows = len(table.rows)
        ppt_cols = len(table.columns)
        excel_rows = len(new_display_texts)
        excel_cols = len(new_display_texts[0]) if new_display_texts else 0

        if ppt_rows != excel_rows or ppt_cols != excel_cols:
            raise ValueError(
                f"表格行列不匹配！PPT 表格为 {ppt_rows}x{ppt_cols}，"
                f"Excel 数据为 {excel_rows}x{excel_cols}（区域: {new_table_value_range}）"
            )

        # 填充格式化后的文本
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                display_text = new_display_texts[i][j]
                if cell.text_frame.paragraphs:
                    p = cell.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = display_text
                        for run in p.runs[1:]:
                            run.text = ""
                    else:
                        p.text = display_text
                else:
                    cell.text = display_text

        return self

    def save(self, new_file_path: str=None):
        new_file_path = new_file_path or self.pptx_file_path
        self.pptx_obj.save(new_file_path)
        return self

    @staticmethod
    def _format_cell_value(cell: Cell) -> str:
        """
        根据 Excel 单元格的 number_format，返回其显示文本（模拟 Excel 显示效果）。
        仅处理常见格式，复杂自定义格式回退为 str(value)。
        """
        value = cell.value
        if value is None:
            return ""

        # 如果是字符串或错误值，直接返回
        if not isinstance(value, (int, float)):
            return str(value)

        nf = cell.number_format.lower().replace("\\", "").replace('"', '').strip()

        # 1. 百分比 %
        if '%' in nf:
            # 计算小数位数：看格式中有几个 0 或 # 在 % 前
            # 简化处理：默认保留2位，或根据格式推断
            if re.search(r'0+\.?0*', nf):
                # 尝试提取小数位数
                decimal_match = re.search(r'\.0+', nf)
                if decimal_match:
                    decimals = len(decimal_match.group().replace('.', ''))
                else:
                    decimals = 0
            else:
                decimals = 2  # 默认
            return f"{value * 100:.{decimals}f}%"

        # 2. 货币（简单处理，只保留数字部分）
        elif any(c in nf for c in ['¥', '$', '€', '£', '元', '人民币']):
            # 提取小数位数
            decimal_match = re.search(r'\.0+', nf)
            decimals = len(decimal_match.group().replace('.', '')) if decimal_match else 2
            # 简单货币符号（这里不处理符号位置，只格式化数字）
            return f"{value:,.{decimals}f}"

        # 3. 数字（带千分位、小数）
        elif re.search(r'[0#,.]+', nf):
            # 检查是否有千分位分隔符
            has_comma = ',' in nf
            # 检查小数位数
            decimal_match = re.search(r'\.0+', nf)
            if decimal_match:
                decimals = len(decimal_match.group().replace('.', ''))
                fmt_str = f"{{:.{decimals}f}}"
                formatted = fmt_str.format(value)
                if has_comma:
                    # 手动加千分位（仅对整数部分）
                    int_part, dec_part = formatted.split('.')
                    int_part = "{:,}".format(int(float(int_part)))
                    formatted = f"{int_part}.{dec_part}"
                return formatted
            else:
                # 无小数
                if has_comma:
                    return "{:,}".format(int(round(value)))
                else:
                    return str(int(round(value)))

        # 4. 常规（General）或未知格式
        else:
            # 尝试模仿 Excel 的 General 格式：去除尾随0
            if isinstance(value, float):
                s = f"{value:.10f}".rstrip('0').rstrip('.')
                return s if s else "0"
            else:
                return str(value)
